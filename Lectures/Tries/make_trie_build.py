from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Alphabet ─────────────────────────────────────────────────────────────────
ALPHA    = ['C','D','E','K','P','R','S','T']
ALPHA_IDX = {c: i for i, c in enumerate(ALPHA)}

# ── Box geometry ─────────────────────────────────────────────────────────────
BW   = 0.230   # letter-box width (inches)
BH   = 0.310   # letter-box height
CKW  = 0.310   # checkmark box width
NW   = len(ALPHA)*BW + CKW   # total node width ≈ 2.15

# ── Colors ───────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D,0x1B,0x2A)
CYAN      = RGBColor(0x00,0xD4,0xFF)
WHITE     = RGBColor(0xE8,0xE8,0xE8)
DIM       = RGBColor(0x66,0x77,0x88)
EMPTY_F   = RGBColor(0x12,0x22,0x32)   # inactive slot
ACTIVE_F  = RGBColor(0x40,0x58,0x72)   # existing active slot
NEW_F     = RGBColor(0xFF,0x99,0x00)   # newly added slot
BORDER    = RGBColor(0x33,0x55,0x77)
WORD_F    = RGBColor(0x00,0xAA,0x55)   # existing end-of-word
NEWWORD_F = RGBColor(0xFF,0xDD,0x00)   # new end-of-word
ARROW_OLD = RGBColor(0x55,0x77,0x99)
ARROW_NEW = RGBColor(0xFF,0x99,0x00)

W, H = 13.33, 7.5

# ── Node top-left positions (x, y) in inches ─────────────────────────────────
# Layout mirrors the image: root/p/pe centered, then branch left (pec*) / right (pet*)
POS = {
    '':       (5.59, 0.28),
    'p':      (5.59, 0.98),
    'pe':     (5.59, 1.68),
    'pec':    (0.70, 2.38),
    'pet':    (8.70, 2.38),
    'peck':   (0.70, 3.08),
    'pete':   (8.00, 3.08),
    'pets':   (10.50, 3.08),
    'pecke':  (0.05, 3.78),
    'pecks':  (2.10, 3.78),
    'peter':  (8.00, 3.78),
    'pecked': (0.05, 4.48),
}

# Parent edge: child_path -> (parent_path, edge_char)
EDGE = {
    'p':      ('',      'P'),
    'pe':     ('p',     'E'),
    'pec':    ('pe',    'C'),
    'pet':    ('pe',    'T'),
    'peck':   ('pec',   'K'),
    'pete':   ('pet',   'E'),
    'pets':   ('pet',   'S'),
    'pecke':  ('peck',  'E'),
    'pecks':  ('peck',  'S'),
    'peter':  ('pete',  'R'),
    'pecked': ('pecke', 'D'),
}

# ── Trie state at each step ───────────────────────────────────────────────────
# State format: path -> {'active': set of chars with children, 'word': bool}
# We track cumulative state + what's new each step.

STEPS = [
    {
        'title':     'Before any insertions — empty root',
        'sub':       'Alphabet: C D E K P R S T',
        'new_nodes': [],
        'new_words': [],
    },
    {
        'title':     'Step 1 — Insert  "pet"',
        'sub':       'Walk P → E → T, creating nodes as needed. Mark T node as end-of-word.',
        'new_nodes': ['p', 'pe', 'pet'],
        'new_words': ['pet'],
    },
    {
        'title':     'Step 2 — Insert  "pets"',
        'sub':       'P → E → T already exist. Add S child to the "pet" node.',
        'new_nodes': ['pets'],
        'new_words': ['pets'],
    },
    {
        'title':     'Step 3 — Insert  "peter"',
        'sub':       'P → E → T already exist. Add E child to "pet", then R child. Mark R node.',
        'new_nodes': ['pete', 'peter'],
        'new_words': ['peter'],
    },
    {
        'title':     'Step 4 — Insert  "peck"',
        'sub':       'P → E already exist. Add C child to "pe" node, then K. Mark K node.',
        'new_nodes': ['pec', 'peck'],
        'new_words': ['peck'],
    },
    {
        'title':     'Step 5 — Insert  "pecked"',
        'sub':       'P → E → C → K already exist. Add E child to "peck", then D. Mark D node.',
        'new_nodes': ['pecke', 'pecked'],
        'new_words': ['pecked'],
    },
    {
        'title':     'Step 6 — Insert  "pecks"',
        'sub':       'P → E → C → K already exist. Add S child to "peck". Mark S node.',
        'new_nodes': ['pecks'],
        'new_words': ['pecks'],
    },
    {
        'title':     'Final Trie',
        'sub':       'pet · pets · peter · peck · pecked · pecks',
        'new_nodes': [],
        'new_words': [],
    },
]

# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = NAVY

def add_text(slide, text, x, y, w, h, size=14, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run()
    r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color; r.font.name = "Calibri"

def draw_node(slide, path, active, is_word, new_active=set(), new_word=False):
    x, y = POS[path]
    for i, ch in enumerate(ALPHA):
        bx = x + i * BW
        rect = slide.shapes.add_shape(
            1, Inches(bx), Inches(y), Inches(BW), Inches(BH))
        rect.line.color.rgb = BORDER
        rect.line.width = Pt(0.75)
        if ch in new_active:
            rect.fill.solid(); rect.fill.fore_color.rgb = NEW_F
            tcol = NAVY
        elif ch in active:
            rect.fill.solid(); rect.fill.fore_color.rgb = ACTIVE_F
            tcol = WHITE
        else:
            rect.fill.solid(); rect.fill.fore_color.rgb = EMPTY_F
            tcol = DIM
        tf = rect.text_frame
        for m in (tf.margin_left, tf.margin_right,
                  tf.margin_top, tf.margin_bottom):
            pass  # can't zero margins via property easily; keep defaults
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ch
        r.font.size = Pt(8); r.font.bold = (ch in active or ch in new_active)
        r.font.color.rgb = tcol; r.font.name = "Courier New"

    # Checkmark box
    bx = x + 8 * BW
    rect = slide.shapes.add_shape(
        1, Inches(bx), Inches(y), Inches(CKW), Inches(BH))
    rect.line.color.rgb = BORDER; rect.line.width = Pt(0.75)
    if new_word:
        rect.fill.solid(); rect.fill.fore_color.rgb = NEWWORD_F
        sym, tcol = "✓", NAVY
    elif is_word:
        rect.fill.solid(); rect.fill.fore_color.rgb = WORD_F
        sym, tcol = "✓", WHITE
    else:
        rect.fill.solid(); rect.fill.fore_color.rgb = EMPTY_F
        sym, tcol = " ", DIM
    tf = rect.text_frame
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r  = p.add_run(); r.text = sym
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = tcol

def char_cx(path, ch):
    """X center of the character box for ch in node at path."""
    x, _ = POS[path]
    return x + ALPHA_IDX[ch] * BW + BW / 2

def draw_arrow(slide, parent_path, child_path, is_new):
    par_ch = EDGE[child_path][1]
    x1 = char_cx(parent_path, par_ch)
    px, py = POS[parent_path]
    y1 = py + BH

    cx, cy = POS[child_path]
    x2 = cx + NW / 2
    y2 = cy

    color = ARROW_NEW if is_new else ARROW_OLD
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.8 if is_new else 1.2)

def draw_label(slide, path, label_text, is_new):
    x, y = POS[path]
    col = NEW_F if is_new else DIM
    add_text(slide, '"' + label_text + '"',
             x + NW + 0.05, y, 0.9, BH+0.02,
             size=9, color=col, italic=True)

# ── Legend helper ─────────────────────────────────────────────────────────────

def draw_legend(slide):
    lx, ly = 10.9, 5.45
    add_text(slide, "Legend", lx, ly, 2.1, 0.35, size=13, bold=True, color=CYAN)
    items = [
        (ACTIVE_F, WHITE, "existing child"),
        (NEW_F,    NAVY,  "new this step"),
        (WORD_F,   WHITE, "word ✓ (existing)"),
        (NEWWORD_F,NAVY,  "word ✓ (new)"),
    ]
    for i, (fill, tc, label) in enumerate(items):
        bx, by = lx, ly + 0.40 + i * 0.38
        rect = slide.shapes.add_shape(1, Inches(bx), Inches(by),
                                      Inches(0.30), Inches(0.26))
        rect.fill.solid(); rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = BORDER; rect.line.width = Pt(0.5)
        add_text(slide, label, bx + 0.36, by - 0.02, 1.75, 0.30,
                 size=11, color=WHITE)

# ── Build cumulative state ────────────────────────────────────────────────────

def build_state(up_to_step):
    """
    Returns (all_nodes, word_nodes, new_nodes_this_step, new_words_this_step).
    all_nodes: set of paths that exist
    word_nodes: set of paths marked as words
    """
    all_nodes  = {''}
    word_nodes = set()
    new_nodes  = set()
    new_words  = set()
    for i, step in enumerate(STEPS):
        if i == 0:
            continue  # step 0 = empty root
        all_nodes.update(step['new_nodes'])
        word_nodes.update(step['new_words'])
        if i == up_to_step:
            new_nodes = set(step['new_nodes'])
            new_words = set(step['new_words'])
        if i >= up_to_step:
            break
    return all_nodes, word_nodes, new_nodes, new_words

def active_chars(path, all_nodes):
    """Return set of chars for which a child of path exists in all_nodes."""
    result = set()
    for child, (par, ch) in EDGE.items():
        if par == path and child in all_nodes:
            result.add(ch)
    return result

# ── Build slides ──────────────────────────────────────────────────────────────

for step_idx, step in enumerate(STEPS):
    s = prs.slides.add_slide(blank)
    set_bg(s)

    # Title bar
    rect = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(W), Inches(0.60))
    rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(0x07,0x14,0x22)
    rect.line.fill.background()
    add_text(s, step['title'], 0.25, 0.04, 9.0, 0.52,
             size=24, bold=True, color=CYAN)
    add_text(s, step['sub'],   0.25, 0.04, 12.8, 0.52,
             size=13, color=DIM, align=PP_ALIGN.RIGHT, italic=True)

    all_nodes, word_nodes, new_nodes, new_words = build_state(step_idx)

    # Draw all edges first (under nodes)
    for child in all_nodes:
        if child == '':
            continue
        par = EDGE[child][0]
        is_new = child in new_nodes
        draw_arrow(s, par, child, is_new)

    # Draw all nodes
    for path in all_nodes:
        ac  = active_chars(path, all_nodes)
        nac = {ch for child, (par, ch) in EDGE.items()
               if par == path and child in new_nodes}
        iw  = path in word_nodes
        nw  = path in new_words
        draw_node(s, path, ac, iw, new_active=nac, new_word=nw)

    # Path label above root
    add_text(s, 'root', POS[''][0] - 0.55, POS[''][1] - 0.01,
             0.52, BH, size=10, color=DIM, align=PP_ALIGN.RIGHT)

    # Word list on the right
    add_text(s, "Words", 11.25, 0.70, 1.85, 0.38, size=14, bold=True, color=CYAN)
    word_list = ['pet','pets','peter','peck','pecked','pecks']
    inserted  = set()
    for i2 in range(1, step_idx + 1):
        inserted.update(STEPS[i2]['new_words'])
    for wi, w in enumerate(word_list):
        col  = (NEW_F if w in new_words else
                WHITE if w in inserted else DIM)
        bold = w in inserted
        add_text(s, ("✓ " if w in inserted else "  ") + w,
                 11.20, 1.08 + wi * 0.50, 1.95, 0.46,
                 size=15, bold=bold, color=col,
                 align=PP_ALIGN.LEFT)

    # Alphabet key
    add_text(s, "Alphabet:", 11.20, 4.25, 1.95, 0.36, size=11, color=DIM)
    add_text(s, " · ".join(ALPHA), 11.20, 4.55, 1.95, 0.36,
             size=11, color=DIM, align=PP_ALIGN.LEFT)

    draw_legend(s)

# ── Save ─────────────────────────────────────────────────────────────────────
out = '/Users/griffin/__currentCourses/5243-Algorithms/Lectures/Tries/trie_build_steps.pptx'
prs.save(out)
print("Saved:", out)
