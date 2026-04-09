from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# Palette
NAVY    = RGBColor(0x0D, 0x1B, 0x2A)
CYAN    = RGBColor(0x00, 0xD4, 0xFF)
GREEN   = RGBColor(0x00, 0xFF, 0x9F)
WHITE   = RGBColor(0xE8, 0xE8, 0xE8)
DIM     = RGBColor(0x77, 0x88, 0x99)
CODE_BG = RGBColor(0x05, 0x0A, 0x14)
RED     = RGBColor(0xFF, 0x55, 0x55)
DKBLUE  = RGBColor(0x05, 0x20, 0x38)

W = 13.33

blank = prs.slide_layouts[6]

# ── helpers ──────────────────────────────────────────────────────────────────

def bg(slide, color=NAVY):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def tb(slide, text, x, y, w, h, size=22, bold=False, italic=False,
       color=WHITE, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name   = font
    return box

def title(slide, text, sub=None):
    tb(slide, text, 0.4, 0.18, W-0.8, 0.95, size=38, bold=True, color=CYAN)
    if sub:
        tb(slide, sub, 0.4, 1.05, W-0.8, 0.55, size=19, color=DIM, italic=True)
    line = slide.shapes.add_shape(
        1, Inches(0.4), Inches(1.55 if sub else 1.10),
        Inches(W-0.8), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

def code_box(slide, code, x, y, w, h, size=16):
    r = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = CODE_BG
    r.line.color.rgb = RGBColor(0x00, 0x44, 0x66)
    box = slide.shapes.add_textbox(
        Inches(x+0.15), Inches(y+0.12), Inches(w-0.3), Inches(h-0.24))
    tf = box.text_frame; tf.word_wrap = False
    for i, ln in enumerate(code.strip().split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(size)
        run.font.name = "Courier New"
        run.font.color.rgb = GREEN

def divider(slide, x, y, h):
    r = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.04), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = DIM
    r.line.fill.background()

# ── Slide 1 — Title ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
tb(s, "Trie", 0.5, 1.4, 12.3, 1.6, size=90, bold=True, color=CYAN,
   align=PP_ALIGN.CENTER)
tb(s, "Not A Tree   (but it is…)", 0.5, 3.0, 12.3, 0.9, size=34,
   italic=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Data Structures & Algorithms", 0.5, 4.1, 12.3, 0.6, size=22,
   color=DIM, align=PP_ALIGN.CENTER)

# ── Slide 2 — What is a Trie ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
title(s, "What is a Trie?", "aka prefix tree  |  aka retrieval tree")
bullets = [
    ("A tree where each node represents a character — not a stored value", WHITE, False),
    ("The root represents the empty string  ε", WHITE, False),
    ("Each edge is labeled with one character", WHITE, False),
    ("A path from root to a marked node ( * ) spells a complete word", WHITE, False),
    ("Nodes are shared — common prefixes travel the same path", CYAN, True),
    ("This shared storage is why prefix queries are so fast", GREEN, False),
]
y = 1.75
for text, col, bld in bullets:
    tb(s, "• " + text, 0.7, y, W-1.4, 0.55, size=21, color=col, bold=bld)
    y += 0.62

# ── Slide 3 — Etymology ──────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
title(s, 'Why "Trie"?')
tb(s, 'From the word:', 0.6, 1.9, 3.0, 0.7, size=28, color=DIM)
tb(s, 're',      3.4,  1.9, 0.8, 0.7, size=34, color=WHITE)
tb(s, 'TRIE',    4.05, 1.9, 1.8, 0.7, size=34, bold=True, color=CYAN)
tb(s, 'val',     5.65, 1.9, 1.5, 0.7, size=34, color=WHITE)
tb(s, 'Officially pronounced  "try"  —  not  "tree"', 0.6, 2.75, W-1.2, 0.6,
   size=26, italic=True, color=DIM)
tb(s, 'Yet it is literally a tree.  That\'s the joke in the title.',
   0.6, 3.45, W-1.2, 0.6, size=24, color=WHITE)
tb(s, '"try"  was chosen to avoid confusion with B-trees, AVL trees, etc.',
   0.6, 4.3, W-1.2, 0.55, size=20, color=DIM)
tb(s, 'The confusion persists anyway.',
   0.6, 4.85, W-1.2, 0.5, size=20, italic=True, color=DIM)

# ── Slide 4 — Visualization ──────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
title(s, "Visualizing a Trie", 'Inserting: "app"  "apt"  "apple"  "apply"')
code_box(s, """          [root]
            |
           (a)
            |
           (p)
          /   \\
        (p)   (t)
       * |     *      <-- "app",  "apt"
        (l)
       /   \\
     (e)   (y)
      *     *         <-- "apple",  "apply" """,
    1.2, 1.7, 7.5, 5.1, size=20)

tb(s, "*  = end of word", 9.1, 2.0, 4.0, 0.55, size=19, color=GREEN)
tb(s, "All four words share\nthe same  a → p  path", 9.1, 2.85, 4.0, 0.9,
   size=19, color=CYAN)
tb(s, "Shared prefix =\nshared nodes =\nspace + time saved",
   9.1, 4.0, 4.0, 1.2, size=19, color=WHITE)

# ── Slide 5 — Why Not a Hash Map ─────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
title(s, "Why Not Just a Hash Map?")
tb(s, "Hash maps are excellent for exact lookups — but prefix queries require scanning every key.",
   0.5, 1.7, W-1.0, 0.55, size=21, color=WHITE)

rows, cols = 5, 4
tbl = s.shapes.add_table(rows, cols,
    Inches(0.5), Inches(2.45), Inches(12.3), Inches(3.8)).table

headers = ["Operation", "Hash Map", "Balanced BST", "Trie"]
data = [
    ["Insert",        "O(k) avg",   "O(k log n)",        "O(k)"],
    ["Exact Search",  "O(k) avg",   "O(k log n)",        "O(k)"],
    ["Prefix Search", "O(k · n)",   "O(k log n + out)",  "O(k + out)"],
    ["Delete",        "O(k) avg",   "O(k log n)",        "O(k)"],
]
for ci, h in enumerate(headers):
    c = tbl.cell(0, ci)
    c.text = h
    c.text_frame.paragraphs[0].font.bold = True
    c.text_frame.paragraphs[0].font.size = Pt(18)
    c.text_frame.paragraphs[0].font.color.rgb = CYAN
    c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    c.fill.solid(); c.fill.fore_color.rgb = DKBLUE

for ri, row in enumerate(data):
    prefix_row = (ri == 2)
    row_bg = RGBColor(0x10,0x28,0x40) if ri%2==0 else RGBColor(0x0D,0x20,0x34)
    for ci, val in enumerate(row):
        c = tbl.cell(ri+1, ci)
        c.text = val
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(17)
        if prefix_row:
            if ci == 1: p.font.color.rgb = RED;   p.font.bold = True
            elif ci == 3: p.font.color.rgb = GREEN; p.font.bold = True
            else: p.font.color.rgb = WHITE
        else:
            p.font.color.rgb = WHITE
        c.fill.solid(); c.fill.fore_color.rgb = row_bg

tb(s, "k = length of string,   n = number of strings stored",
   0.5, 6.42, W-1.0, 0.4, size=15, color=DIM, italic=True)

# ── Slide 6 — Applications Overview ─────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
title(s, "Real-World Applications")
apps = [
    (CYAN,  "Autocomplete / Suggestions",
             "IDE symbol completion, phone contacts, search bars"),
    (GREEN, "IP Routing",
             "Longest prefix match — runs in every router on the internet"),
    (WHITE, "Spell Check",
             "Detection O(k);  suggestions via edit distance enumeration"),
    (DIM,   "DNS Resolution",
             "Domain hierarchy (com → google → mail) is a natural trie"),
    (DIM,   "Git & Filesystem",
             "Tab-completion of branch names and file paths"),
]
y = 1.85
for col, head, desc in apps:
    tb(s, head, 0.7, y,       W-1.4, 0.48, size=22, bold=True,  color=col)
    tb(s, desc, 0.7, y+0.46, W-1.4, 0.38, size=18, color=DIM)
    y += 1.0

# ── Slide 7 — Autocomplete vs Suggestions ───────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
title(s, "Autocomplete vs. Suggestions", "They are solving different problems")

tb(s, "Autocomplete", 0.5, 1.8, 6.0, 0.55, size=26, bold=True, color=CYAN)
for i, line in enumerate([
    "One right answer",
    'IDE:  std::vec  →  std::vector',
    "Terminal:  git che<TAB>  →  git checkout",
    "Trie walks prefix, returns subtree",
    "Pick first (or only) match.  Done.",
]):
    tb(s, "• " + line, 0.7, 2.45+i*0.54, 5.9, 0.5, size=19, color=WHITE)

divider(s, 6.65, 1.8, 5.1)

tb(s, "Suggestions  (Google)", 6.9, 1.8, 6.0, 0.55, size=26, bold=True, color=GREEN)
for i, line in enumerate([
    "Many candidates, must be ranked",
    'Type "algo" → Trie returns millions',
    "Rank by: global frequency, your",
    "history, location, trending, ML…",
    "Trie does the easy part.",
]):
    tb(s, "• " + line, 7.1, 2.45+i*0.54, 5.9, 0.5, size=19, color=WHITE)

tb(s, 'The Trie handles "find the candidates."  Everything after that is outside the Trie\'s job.',
   0.5, 6.7, W-1.0, 0.5, size=16, color=DIM, italic=True)

# ── Slide 8 — IP Routing ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
title(s, "IP Routing — Longest Prefix Match")
tb(s, "Every IP address is 32 bits (IPv4). Routing entries are CIDR blocks: 192.168.1.0/24",
   0.5, 1.7, W-1.0, 0.5, size=20, color=WHITE)

code_box(s, """/24  →  first 24 bits must match, last 8 can be anything

Routing table (binary Trie, one node per bit):
  10.0.0.0/8       match first  8 bits  →  mark here
  10.20.0.0/16     match first 16 bits  →  mark here
  10.20.30.0/24    match first 24 bits  →  mark here

For each incoming packet:
  Walk the Trie bit-by-bit.
  Track the LAST node that had a valid route entry.
  That deepest match wins.  (Longest Prefix Match)""",
    0.5, 2.35, 8.5, 4.35, size=17)

tb(s, "Cost:",              9.3, 2.4,  3.7, 0.45, size=20, bold=True, color=CYAN)
tb(s, "O(32) — IPv4\nO(128) — IPv6", 9.3, 2.9, 3.7, 0.85, size=21, bold=True, color=GREEN)
tb(s, "Independent of\nrouting table size", 9.3, 3.9, 3.7, 0.75, size=19, color=WHITE)
tb(s, "Without Trie:\nO(n) scan per packet\n→ unacceptable at\nline rate",
   9.3, 4.8, 3.7, 1.4, size=18, color=RED)

# ── Slide 9 — Spell Check ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
title(s, "Spell Check — Two Distinct Problems")

tb(s, "Phase 1 — Detection",  0.5, 1.75, 6.2, 0.5, size=24, bold=True, color=CYAN)
tb(s, "Is the word in the dictionary?", 0.7, 2.28, 6.0, 0.4, size=20, color=WHITE)
tb(s, "Trie exact-match lookup:  O(k).  Done.", 0.7, 2.70, 6.0, 0.4, size=20, color=GREEN, bold=True)

tb(s, "Phase 2 — Suggestion",  0.5, 3.3, 6.2, 0.5, size=24, bold=True, color=CYAN)
for i, line in enumerate([
    "Word not found — what did they mean?",
    "Metric: Levenshtein distance",
    '  (insertions, deletions, substitutions)',
    '"recieve" → "receive"  distance = 1',
    "Generate all strings within dist 1 of W",
    "Look each up in the Trie",
    "Return the ones that exist",
]):
    tb(s, line, 0.7, 3.85+i*0.43, 6.2, 0.4, size=18, color=WHITE)

divider(s, 7.1, 1.75, 5.3)

tb(s, "Complexity",           7.4, 1.75, 5.8, 0.5, size=22, bold=True, color=CYAN)
tb(s, "dist-1 candidates\nfrom a length-k word:", 7.4, 2.35, 5.8, 0.7, size=19, color=DIM)
tb(s, "O(54k) candidates",    7.4, 3.1,  5.8, 0.5, size=22, bold=True, color=GREEN)
tb(s, "× O(k) per Trie lookup", 7.4, 3.65, 5.8, 0.45, size=19, color=WHITE)
tb(s, "= O(k²) total",        7.4, 4.15, 5.8, 0.5, size=26, bold=True, color=CYAN)
tb(s, "Fast enough for\nevery keystroke",  7.4, 4.85, 5.8, 0.7, size=19, color=WHITE)

# ── Slide 10 — C++ Structures ────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
title(s, "C++ — Two Common Node Designs")

code_box(s, """// Design A: Self-referential
class Trie {
private:
    unordered_map<char,Trie*> children;
    bool is_end = false;
public:
    Trie() {}
    void insert(string word);
    vector<string> search(string prefix);
};""", 0.4, 1.7, 6.15, 4.0, size=17)

code_box(s, """// Design B: Separate node class
class TrieNode {
public:
    unordered_map<char,TrieNode*> children;
    bool is_word = false;
};
class Trie {
    TrieNode* root;
public:
    Trie() { root = new TrieNode(); }
    void insert(string word);
    bool search(string word);
};""", 6.75, 1.7, 6.15, 4.6, size=17)

tb(s, "Both use  unordered_map<char, Node*>  →  O(1) average child lookup  |  Max 26 or 128 children per node",
   0.4, 6.4, W-0.8, 0.5, size=16, color=DIM, italic=True)

# ── Slide 11 — Insert & Search ───────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
title(s, "Insert & Prefix Search")

code_box(s, """void insert(string word) {
    Trie* node = this;
    for (char c : word) {
        if (!node->children.count(c))
            node->children[c] = new Trie();
        node = node->children[c];
    }
    node->is_end = true; // mark word end
}""", 0.4, 1.7, 6.15, 3.6, size=17)

code_box(s, """vector<string> search(string prefix) {
    Trie* node = this;
    for (char c : prefix) {
        if (!node->children.count(c))
            return {};     // prefix absent
        node = node->children[c];
    }
    return find_words(node, prefix);
}
// find_words: DFS from here,
// collect every is_end node below""", 6.75, 1.7, 6.15, 4.4, size=17)

tb(s, "Insert: O(k)     |     Prefix search: O(k + output size)",
   0.4, 6.5, W-0.8, 0.5, size=20, bold=True, color=CYAN)

# ── Slide 12 — Summary ───────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
title(s, "Summary")
rows = [
    ("Structure",    "Tree where paths spell words. Common prefixes share nodes."),
    ("Insert",       "Walk one edge per character. O(k). Create nodes as needed."),
    ("Exact search", "Walk edges. If you fall off, it's not there. O(k)."),
    ("Prefix search","Walk prefix, collect subtree. O(k + output). Hash map can't do this cheaply."),
    ("vs Hash Map",  "Prefix query: O(k·n) hash map vs O(k + out) Trie. Not close."),
    ("Real uses",    "Autocomplete, IP routing (TCAM), spell check, git tab-complete."),
    ("Key insight",  "Shared prefixes = shared nodes = fast prefix queries."),
]
y = 1.75
for label, desc in rows:
    tb(s, label + ":", 0.55, y, 2.65, 0.48, size=19, bold=True, color=CYAN)
    tb(s, desc,        3.2,  y, 9.8,  0.48, size=19, color=WHITE)
    y += 0.66

# ── Save ─────────────────────────────────────────────────────────────────────
out = '/Users/griffin/__currentCourses/5243-Algorithms/Lectures/Tries/tries_lecture.pptx'
prs.save(out)
print("Saved:", out)
